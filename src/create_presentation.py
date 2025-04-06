from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
import os
import glob

def add_title_slide(prs, title, subtitle):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

def add_content_slide(prs, title, content):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = title
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

def add_image_slide(prs, title, image_path):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Add image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), height=Inches(5.5))

def add_two_images_slide(prs, title, image_path1, image_path2):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Add images side by side
    if os.path.exists(image_path1):
        slide.shapes.add_picture(image_path1, Inches(0.5), Inches(1.5), height=Inches(5))
    if os.path.exists(image_path2):
        slide.shapes.add_picture(image_path2, Inches(6), Inches(1.5), height=Inches(5))

def create_presentation():
    # Create results directory if it doesn't exist
    os.makedirs('results', exist_ok=True)
    
    prs = Presentation()
    
    # Title Slide
    add_title_slide(prs, 
                   "Comparative Analysis of A* and UCS Algorithms",
                   "Graph Theory Project")

    # Introduction
    add_content_slide(prs, "Introduction", [
        "• Comparison of A* (A-star) and UCS (Uniform Cost Search) algorithms",
        "• Performance analysis on different graph sizes",
        "• Impact of heuristic function",
        "• Evaluation based on expanded nodes and execution time",
        "• Visualization of A* paths for different graph sizes"
    ])

    # Methodology
    add_content_slide(prs, "Methodology", [
        "• Tested graph sizes:",
        "  - 7 different levels (from 5 to 35 vertices)",
        "  - 10 random graphs per level",
        "• Measured metrics:",
        "  - Number of expanded nodes",
        "  - Algorithm execution time",
        "  - Improvement percentage",
        "• Heuristic function: Minimum weighted distance",
        "• Path visualization for sample graphs"
    ])

    # Implementation Details
    add_content_slide(prs, "Implementation Details", [
        "• Graph Structure:",
        "  - Adjacency list representation",
        "  - Weighted edges (range 1-20)",
        "  - Connected graph guarantee",
        "• A* Algorithm:",
        "  - f(n) = g(n) + h(n)",
        "  - Admissible heuristic",
        "• UCS Algorithm:",
        "  - Dijkstra-based implementation",
        "  - Using only g(n)"
    ])

    try:
        # Read experimental data
        df = pd.read_csv('src/algorithm_comparison.csv')
        
        # Add overall comparison
        add_image_slide(prs, "Overall Performance Comparison", 
                       "results/overall_comparison.png")

        # Get all A* path images
        astar_paths = sorted(glob.glob("src/AStar_path_*.png"))
        
        # Add A* Path Visualization Section
        add_content_slide(prs, "A* Path Visualizations", [
            "• Sample optimal paths found by A* algorithm",
            "• Different graph sizes demonstrate algorithm's adaptability",
            "• Path colors indicate:",
            "  - Red: Start node",
            "  - Green: Goal node",
            "  - Blue: Optimal path",
            "• Edge weights shown on connections"
        ])
        
        # Add A* path visualizations
        for path_image in astar_paths:
            add_image_slide(prs, 
                          f"A* Path Visualization - Sample Graph",
                          path_image)

        # Results for each graph size
        for vertices in sorted(df['vertices'].unique()):
            level_data = df[df['vertices'] == vertices]
            edges = level_data['edges'].iloc[0]
            
            # Add performance comparison
            add_image_slide(prs, 
                          f"Performance Analysis (V={vertices}, E={edges})",
                          f"results/comparison_V{vertices}_E{edges}.png")

        # Key Findings
        improvement_by_size = df.groupby('vertices').agg({
            'expanded_nodes_astar': 'mean',
            'expanded_nodes_ucs': 'mean'
        }).reset_index()
        
        improvement_percentage = ((improvement_by_size['expanded_nodes_ucs'] - 
                                 improvement_by_size['expanded_nodes_astar']) / 
                                 improvement_by_size['expanded_nodes_ucs'] * 100)
        
        max_improvement = improvement_percentage.max()
        avg_improvement = improvement_percentage.mean()
        
        add_content_slide(prs, "Key Findings", [
            f"• Average improvement: {avg_improvement:.1f}%",
            f"• Maximum improvement: {max_improvement:.1f}%",
            "• A* algorithm advantages:",
            "  - Fewer node expansions",
            "  - Increasing performance gap with graph size",
            "  - Consistent execution time",
            "• Impact of heuristic function:",
            "  - Significant reduction in search space",
            "  - Maintaining optimal solution guarantee",
            "• Visual confirmation of optimal paths"
        ])

        # Detailed Analysis
        add_content_slide(prs, "Detailed Analysis", [
            "• Small Graphs (V=5):",
            "  - Minimal difference between A* and UCS",
            "  - Both algorithms efficient",
            "  - Clear path visualization",
            "• Medium-sized Graphs (V=15-25):",
            "  - A* advantage becomes apparent",
            "  - UCS expanded nodes increase significantly",
            "  - More complex path structures",
            "• Large Graphs (V=30-35):",
            "  - A* maintains consistent performance",
            "  - UCS performance degrades significantly",
            "  - Sophisticated path finding demonstrated"
        ])

        # Conclusion
        add_content_slide(prs, "Conclusion", [
            "• A* algorithm is more effective with appropriate heuristic function",
            "• A* advantage increases with graph size",
            "• Heuristic function significantly reduces search space",
            "• Optimal path guarantee is maintained",
            "• Visual examples confirm efficient path finding",
            "• A* is recommended for practical applications"
        ])

    except FileNotFoundError:
        print("Error: 'algorithm_comparison.csv' file not found!")
        print("Please run the C++ code first to generate results.")
        return
    except Exception as e:
        print(f"Unexpected error occurred: {str(e)}")
        return

    try:
        # Save presentation
        prs.save('results/algorithm_comparison_presentation.pptx')
        print("Presentation created successfully!")
        print("File location: results/algorithm_comparison_presentation.pptx")
    except Exception as e:
        print(f"Error while saving presentation: {str(e)}")

if __name__ == "__main__":
    create_presentation() 